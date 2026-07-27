from base64 import b64decode
from io import BytesIO
import re

from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def generate_pptx_file(title: str, image_b64: str):
    image_b64_clean = re.sub(r"^data:image/.+;base64,", "", image_b64)
    image_bytes = b64decode(image_b64_clean)

    image = Image.open(BytesIO(image_bytes)).convert("RGB")
    img_width_px, img_height_px = image.size

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    img_stream = BytesIO()
    image.save(img_stream, format="PNG")
    img_stream.seek(0)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    dpi = 96.0
    img_w_in = img_width_px / dpi
    img_h_in = img_height_px / dpi

    max_w_in = slide_w / 914400.0
    max_h_in = slide_h / 914400.0

    scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    final_w = img_w_in * scale
    final_h = img_h_in * scale

    left = (max_w_in - final_w) / 2.0
    top = (max_h_in - final_h) / 2.0

    slide.shapes.add_picture(
        img_stream,
        Inches(left),
        Inches(top),
        width=Inches(final_w),
        height=Inches(final_h),
    )

    out = BytesIO()
    prs.save(out)
    out.seek(0)

    safe_title = re.sub(r"[^a-zA-Z0-9 _\-\.\(\)]", "", title).strip() or "LifeViz"
    return out, safe_title
