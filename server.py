from flask import Flask, request, send_file, jsonify, render_template
from io import BytesIO
from base64 import b64decode
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import pandas as pd
import re
import json
import os

from openai import OpenAI

# OpenAI client reads OPENAI_API_KEY from your environment
client = OpenAI()

app = Flask(__name__, static_folder='static', template_folder='templates')


# ===================== Utility helpers =====================

def _normalize_col(name: str) -> str:
    """Lower-case, collapse spaces for matching column names case-insensitively."""
    return re.sub(r'\s+', ' ', str(name)).strip().lower()


def _pretty_label(col_name: str) -> str:
    """Create a simple human-friendly label from a column name."""
    return re.sub(r'\s+', ' ', str(col_name)).strip().capitalize()


def _clean_cell(val):
    if pd.isna(val):
        return ""
    text = str(val).strip()
    return "" if not text or text.lower() == 'nan' else text


def _excelish_month(val):
    """
    Convert workbook timing values into month labels the frontend can parse.
    Supports Excel serial dates, timestamps, and month-like strings.
    """
    if pd.isna(val):
        return ""

    if isinstance(val, (int, float)) and not isinstance(val, bool):
        try:
            dt = pd.to_datetime("1899-12-30") + pd.to_timedelta(float(val), unit="D")
            return dt.strftime("%b")
        except Exception:
            return str(val).strip()

    try:
        dt = pd.to_datetime(val, errors="coerce")
        if pd.notna(dt):
            return dt.strftime("%b")
    except Exception:
        pass

    text = str(val).strip()
    if not text:
        return ""

    lower = text.lower()
    month_map = {
        "jan": "Jan", "january": "Jan",
        "feb": "Feb", "february": "Feb",
        "mar": "Mar", "march": "Mar",
        "apr": "Apr", "april": "Apr",
        "may": "May",
        "jun": "Jun", "june": "Jun",
        "jul": "Jul", "july": "Jul",
        "aug": "Aug", "august": "Aug",
        "sep": "Sep", "sept": "Sep", "september": "Sep",
        "oct": "Oct", "october": "Oct",
        "nov": "Nov", "november": "Nov",
        "dec": "Dec", "december": "Dec",
    }
    for key, label in month_map.items():
        if key in lower:
            return label
    return text


def _range_text_from_values(values, suffix=""):
    cleaned = [str(v).strip() for v in values if str(v).strip()]
    if not cleaned:
        return ""
    uniq = list(dict.fromkeys(cleaned))
    if len(uniq) == 1:
        return f"{uniq[0]}{suffix}".strip()
    return f"{uniq[0]}-{uniq[-1]}{suffix}".strip()


def _sheet_name_matches(sheet_name: str, *needles: str) -> bool:
    normalized = _normalize_col(sheet_name)
    return any(needle in normalized for needle in needles)


def _sheet_cell(df, row_idx, col_idx):
    if row_idx >= len(df.index) or col_idx >= len(df.columns):
        return ""
    return _clean_cell(df.iat[row_idx, col_idx])


def _coalesce_label_value(items, max_items=4):
    cleaned = []
    for label, value in items:
        label_text = _clean_cell(label)
        value_text = _clean_cell(value)
        if not label_text or not value_text:
            continue
        cleaned.append(f"{label_text}: {value_text}")
    return cleaned[:max_items]


def _lifespan_unit_suffix(metric, value=""):
    unit = _clean_cell(metric).lower()
    value_text = _clean_cell(value)
    singular = value_text in ("1", "1.0")
    if "year" in unit:
        return " year" if singular else " years"
    if "month" in unit:
        return " month" if singular else " months"
    if "week" in unit:
        return " week" if singular else " weeks"
    if "day" in unit:
        return " day" if singular else " days"
    return f" {metric}".rstrip() if metric else ""


def _select_lifespan_text(lifespan_rows):
    """
    Prefer one representative lifespan value, not a min/max range.
    Priority: exact mean/average, average-range maximum, average-range minimum,
    maximum, then the first usable row.
    """
    cleaned = []
    for row in lifespan_rows:
        if len(row) == 3:
            metric, value, qualifier = row
        else:
            metric, value = row
            qualifier = ""
        metric = _clean_cell(metric)
        value = _clean_cell(value)
        qualifier = _clean_cell(qualifier)
        if value:
            cleaned.append((metric, value, qualifier))

    if not cleaned:
        return ""

    def pick(predicate):
        return next((row for row in cleaned if predicate(row[2].lower())), None)

    selected = (
        pick(lambda q: ("mean" in q or "average" in q) and "range" not in q and "min" not in q and "max" not in q)
        or pick(lambda q: "average range" in q and "max" in q)
        or pick(lambda q: "average range" in q and "min" in q)
        or pick(lambda q: q.strip() == "maximum" or q.strip() == "max")
        or cleaned[0]
    )

    metric, value, _ = selected
    return f"{value}{_lifespan_unit_suffix(metric, value)}".strip()


def _structured_life_history_workbook_to_species_list(raw_sheets):
    """
    Parse the worksheet format used in the attached RISE examples.
    Works when the user uploads either:
    - several worksheet files that together form one species workbook, or
    - a single life-history worksheet file.
    """
    stage_sheet_name = next(
        (name for name in raw_sheets if _sheet_name_matches(name, "life history stages", "lifestages")),
        None
    )
    if not stage_sheet_name:
        return None

    species_map = {}

    def get_species_entry(scientific_name, common_name="", population=""):
        sci = _clean_cell(scientific_name)
        if not sci:
            return None
        key = sci.lower()
        entry = species_map.setdefault(key, {
            "name": sci,
            "commonName": _clean_cell(common_name),
            "population": _clean_cell(population),
            "imageFile": None,
            "imageUrl": None,
            "lifespan": [],
            "stages": {}
        })
        if common_name and not entry["commonName"]:
            entry["commonName"] = _clean_cell(common_name)
        if population and not entry["population"]:
            entry["population"] = _clean_cell(population)
        return entry

    def get_stage_entry(species_entry, stage_name):
        stage = _clean_cell(stage_name)
        if not species_entry or not stage:
            return None
        key = stage.lower()
        return species_entry["stages"].setdefault(key, {
            "name": stage,
            "duration_values": [],
            "duration_units": [],
            "timing_months": [],
            "bullets": []
        })

    stage_df = raw_sheets[stage_sheet_name]
    for row_idx in range(2, len(stage_df.index)):
        species_entry = get_species_entry(
            _sheet_cell(stage_df, row_idx, 0),
            _sheet_cell(stage_df, row_idx, 1),
            _sheet_cell(stage_df, row_idx, 2)
        )
        stage_entry = get_stage_entry(species_entry, _sheet_cell(stage_df, row_idx, 3))
        if not stage_entry:
            continue

        duration_metric = _sheet_cell(stage_df, row_idx, 6)
        duration_value = _sheet_cell(stage_df, row_idx, 7)
        if duration_metric and duration_value:
            stage_entry["duration_units"].append(duration_metric)
            stage_entry["duration_values"].append(duration_value)

        timing_value = _sheet_cell(stage_df, row_idx, 10)
        if timing_value:
            stage_entry["timing_months"].append(_excelish_month(timing_value))

        reproductive_age_value = _sheet_cell(stage_df, row_idx, 16)
        reproductive_age_metric = _sheet_cell(stage_df, row_idx, 15)
        if reproductive_age_value:
            suffix = f" {reproductive_age_metric}".strip()
            stage_entry["bullets"].append(f"Reproduction age: {reproductive_age_value}{suffix}")

        reproductive_frequency_value = _sheet_cell(stage_df, row_idx, 19)
        reproductive_frequency_metric = _sheet_cell(stage_df, row_idx, 18)
        if reproductive_frequency_value:
            suffix = f" {reproductive_frequency_metric}".strip()
            stage_entry["bullets"].append(f"Reproductive frequency: {reproductive_frequency_value}{suffix}")

        initial_count_value = _sheet_cell(stage_df, row_idx, 22)
        initial_count_metric = _sheet_cell(stage_df, row_idx, 21)
        if initial_count_value:
            suffix = f" {initial_count_metric}".strip()
            stage_entry["bullets"].append(f"Initial count: {initial_count_value}{suffix}")

        repro_strategy = _sheet_cell(stage_df, row_idx, 23)
        if repro_strategy:
            stage_entry["bullets"].append(f"Reproduction: {repro_strategy}")

    habitat_sheet_name = next((name for name in raw_sheets if _sheet_name_matches(name, "habitat")), None)
    if habitat_sheet_name:
        habitat_df = raw_sheets[habitat_sheet_name]
        habitat_names = [_sheet_cell(habitat_df, 1, col_idx) for col_idx in range(4, len(habitat_df.columns))]

        for row_idx in range(2, len(habitat_df.index)):
            species_entry = get_species_entry(
                _sheet_cell(habitat_df, row_idx, 0),
                _sheet_cell(habitat_df, row_idx, 1),
                _sheet_cell(habitat_df, row_idx, 2)
            )
            stage_entry = get_stage_entry(species_entry, _sheet_cell(habitat_df, row_idx, 3))
            if not stage_entry:
                continue

            active = []
            for offset, habitat_name in enumerate(habitat_names, start=4):
                marker = _sheet_cell(habitat_df, row_idx, offset).lower()
                if marker in ("x", "yes", "true", "1"):
                    active.append(habitat_name)
            if active:
                stage_entry["bullets"].append(f"Habitat: {', '.join(active)}")

    resource_sheet_name = next((name for name in raw_sheets if _sheet_name_matches(name, "resource needs")), None)
    if resource_sheet_name:
        resource_df = raw_sheets[resource_sheet_name]
        need_names = [_sheet_cell(resource_df, 1, col_idx) for col_idx in range(4, len(resource_df.columns))]

        for row_idx in range(2, len(resource_df.index)):
            species_entry = get_species_entry(
                _sheet_cell(resource_df, row_idx, 0),
                _sheet_cell(resource_df, row_idx, 1),
                _sheet_cell(resource_df, row_idx, 2)
            )
            stage_entry = get_stage_entry(species_entry, _sheet_cell(resource_df, row_idx, 3))
            if not stage_entry:
                continue

            resource_pairs = []
            for offset, need_name in enumerate(need_names, start=4):
                resource_value = _sheet_cell(resource_df, row_idx, offset)
                if resource_value:
                    resource_pairs.append((need_name, resource_value))
            if resource_pairs:
                stage_entry["bullets"].append(
                    f"Resource needs: {'; '.join(_coalesce_label_value(resource_pairs, 4))}"
                )

    lifespan_sheet_name = next((name for name in raw_sheets if _sheet_name_matches(name, "life span", "lifespan")), None)
    if lifespan_sheet_name:
        lifespan_df = raw_sheets[lifespan_sheet_name]
        for row_idx in range(1, len(lifespan_df.index)):
            species_entry = get_species_entry(
                _sheet_cell(lifespan_df, row_idx, 0),
                _sheet_cell(lifespan_df, row_idx, 1),
                ""
            )
            if not species_entry:
                continue

            qualifier = _sheet_cell(lifespan_df, row_idx, 3)
            metric = _sheet_cell(lifespan_df, row_idx, 4)
            value = _sheet_cell(lifespan_df, row_idx, 5)
            if metric and value:
                species_entry["lifespan"].append((metric, value, qualifier))

    species_out = []
    for info in species_map.values():
        title_base = info["commonName"] or info["name"]
        title = f"{title_base} - Life history"
        if info["population"]:
            title = f"{title_base} ({info['population']}) - Life history"

        lifespan_text = _select_lifespan_text(info["lifespan"])

        stages_out = []
        for st in info["stages"].values():
            bullets = []

            unit_suffix = f" {st['duration_units'][0]}" if st["duration_units"] else ""
            duration_text = _range_text_from_values(st["duration_values"], unit_suffix)
            if duration_text:
                bullets.append(f"Duration: {duration_text}")

            months = [m for m in st["timing_months"] if m]
            if months:
                bullets.append(f"Timing: {_range_text_from_values(months)}")

            if lifespan_text:
                bullets.append(f"Lifespan: {lifespan_text}")

            bullets.extend(st["bullets"])
            stages_out.append({
                "title": st["name"],
                "bullets": list(dict.fromkeys([b for b in bullets if b]))
            })

        species_out.append({
            "name": info["name"],
            "title": title,
            "imageFile": info["imageFile"],
            "imageUrl": info["imageUrl"],
            "lifespan": lifespan_text,
            "stages": stages_out
        })

    return species_out or None


def _generic_workbook_to_species_list(sheets):
    species_dict = {}

    for _, df in sheets.items():
        if df is None or df.empty:
            continue

        normalized_cols = {_normalize_col(c): c for c in df.columns}

        species_col = None
        for key, orig in normalized_cols.items():
            if 'species' in key:
                species_col = orig
                break
        if not species_col:
            continue

        stage_col = None
        for key, orig in normalized_cols.items():
            if 'stage' in key or 'lifestage' in key or 'life stage' in key:
                stage_col = orig
                break

        for _, row in df.iterrows():
            _add_bullets_from_row(df, row, species_dict, species_col, stage_col)

    if not species_dict:
        return None

    species_list = []
    for info in species_dict.values():
        stages = []

        for st in info["stages"].values():
            bullets = []
            seen = set()
            for bullet in st["bullets"]:
                if bullet not in seen:
                    bullets.append(bullet)
                    seen.add(bullet)
            stages.append({"title": st["name"], "bullets": bullets})

        if info["speciesBullets"]:
            stages.insert(0, {
                "title": "Species summary",
                "bullets": list(dict.fromkeys(info["speciesBullets"]))
            })

        species_list.append({
            "name": info["name"],
            "title": f"{info['name']} - Life history",
            "imageFile": info.get("imageFile"),
            "imageUrl": info.get("imageUrl"),
            "stages": stages
        })

    return species_list


def _simple_template_sheet_to_species_list(df):
    if df is None or df.empty:
        return None

    cols_norm = {_normalize_col(c): c for c in df.columns}

    def col(*names):
        for n in names:
            key = _normalize_col(n)
            if key in cols_norm:
                return cols_norm[key]
        return None

    species_col = col('Species')
    stage_col = col('StageName', 'Stage', 'Life stage', 'Lifestage')
    order_col = col('StageOrder', 'Order')
    common_col = col('CommonName', 'Common name')
    image_file_col = col('ImageFile', 'Image', 'Photo')
    image_url_col = col('ImageUrl', 'Image URL', 'PhotoUrl', 'Photo URL')

    if not species_col or not stage_col:
        return None

    detail_candidates = [
        'Duration', 'Timing', 'Habitat', 'Lifespan', 'Movement',
        'Physical', 'Reproduction', 'Food', 'Notes'
    ]
    detail_cols = [col(c) for c in detail_candidates if col(c)]

    species_map = {}

    for _, row in df.iterrows():
        raw_species = row[species_col]
        if pd.isna(raw_species):
            continue
        s_name = str(raw_species).strip()
        if not s_name or s_name.lower() == 'nan':
            continue

        s_key = s_name.lower()
        sp = species_map.setdefault(
            s_key,
            {
                "name": s_name,
                "commonName": None,
                "imageFile": None,
                "imageUrl": None,
                "stages": {}
            }
        )

        if common_col and not pd.isna(row[common_col]):
            sp["commonName"] = str(row[common_col]).strip()
        if image_file_col and not pd.isna(row[image_file_col]):
            sp["imageFile"] = str(row[image_file_col]).strip()
        if image_url_col and not pd.isna(row[image_url_col]):
            sp["imageUrl"] = str(row[image_url_col]).strip()

        raw_stage = row[stage_col]
        if pd.isna(raw_stage):
            continue
        stage_name = str(raw_stage).strip()
        if not stage_name or stage_name.lower() == 'nan':
            continue

        order = None
        if order_col and not pd.isna(row[order_col]):
            try:
                order = float(row[order_col])
            except Exception:
                order = None

        st_key = stage_name.lower()
        stage = sp["stages"].setdefault(
            st_key,
            {"name": stage_name, "order": order, "bullets": []}
        )

        for detail_col in detail_cols:
            val = row[detail_col]
            if pd.isna(val):
                continue
            text = str(val).strip()
            if not text or text.lower() == 'nan':
                continue
            stage["bullets"].append(f"{_pretty_label(detail_col)}: {text}")

    if not species_map:
        return None

    species_out = []
    for sp in species_map.values():
        stages = list(sp["stages"].values())
        stages.sort(key=lambda st: (9999 if st["order"] is None else st["order"]))
        species_out.append({
            "name": sp["name"],
            "title": f"{sp['name']} - Life history",
            "imageFile": sp["imageFile"],
            "imageUrl": sp["imageUrl"],
            "stages": [{"title": st["name"], "bullets": st["bullets"]} for st in stages]
        })

    return species_out


def _read_excel_sheets_from_bytes(file_bytes, header=0):
    return pd.read_excel(BytesIO(file_bytes), sheet_name=None, header=header)


def _merge_sheet_maps(sheet_maps):
    merged = {}
    for sheets in sheet_maps:
        for name, df in sheets.items():
            unique_name = name
            suffix = 2
            while unique_name in merged:
                unique_name = f"{name} ({suffix})"
                suffix += 1
            merged[unique_name] = df
    return merged


def _dataset_response(title, species_list):
    return {"title": title, "species": species_list}


def _add_bullets_from_row(df, row, species_dict, species_col, stage_col=None):
    """
    Generic aggregator for the full workbook route.

    - Uses species + optional stage columns (case-insensitive) to group rows.
    - All other non-empty columns become bullets like "Label: value".
    - Species and stage keys are lowercased so differing capitalization merges.
    - Display names keep the first capitalization encountered.
    """
    raw_species = row[species_col]
    if pd.isna(raw_species):
        return

    s_name = str(raw_species).strip()
    if not s_name or s_name.lower() == 'nan':
        return

    s_key = s_name.lower()

    species_entry = species_dict.setdefault(
        s_key,
        {
            "key": s_key,
            "name": s_name,
            "commonName": None,
            "imageFile": None,
            "imageUrl": None,
            "stages": {},
            "speciesBullets": [],
        }
    )

    stage_name = None
    if stage_col is not None and stage_col in df.columns:
        raw_stage = row[stage_col]
        if not pd.isna(raw_stage):
            tmp = str(raw_stage).strip()
            if tmp and tmp.lower() != 'nan':
                stage_name = tmp

    if stage_name:
        st_key = stage_name.lower()
        stage = species_entry["stages"].setdefault(
            st_key,
            {"key": st_key, "name": stage_name, "bullets": []}
        )
        target_list = stage["bullets"]
    else:
        target_list = species_entry["speciesBullets"]

    for col in df.columns:
        if col == species_col or col == stage_col:
            continue
        val = row[col]
        if pd.isna(val):
            continue
        text = str(val).strip()
        if not text or text.lower() == 'nan':
            continue
        label = _pretty_label(col)
        bullet = f"{label}: {text}"
        target_list.append(bullet)


# ===================== Routes =====================

@app.route('/')
def index():
    return render_template('index.html')


# ---------- Export PPTX from PNG ----------

@app.route('/generate_pptx', methods=['POST'])
def generate_pptx():
    data = request.get_json()
    if not data or 'image_b64' not in data:
        return jsonify({"error": "Missing image_b64"}), 400

    title = data.get('title', 'Export')
    image_b64_clean = re.sub(r'^data:image/.+;base64,', '', data['image_b64'])
    image_bytes = b64decode(image_b64_clean)

    image = Image.open(BytesIO(image_bytes)).convert('RGB')
    img_width_px, img_height_px = image.size

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

    img_stream = BytesIO()
    image.save(img_stream, format='PNG')
    img_stream.seek(0)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    dpi = 96.0
    img_w_in = img_width_px / dpi
    img_h_in = img_height_px / dpi

    max_w_in = slide_w / 914400.0
    max_h_in = slide_h / 914400.0

    scale = min(max_w_in / img_w_in, max_h_in / img_h_in, 1.0)
    final_w = img_w_in * scale
    final_h = img_h_in * scale

    left = (max_w_in - final_w) / 2.0
    top = (max_h_in - final_h) / 2.0

    slide.shapes.add_picture(
        img_stream,
        Inches(left),
        Inches(top),
        width=Inches(final_w),
        height=Inches(final_h),
    )

    out = BytesIO()
    prs.save(out)
    out.seek(0)

    safe_title = re.sub(r'[^a-zA-Z0-9 _\-\.\(\)]', '', title).strip() or "LifeViz"
    return send_file(
        out,
        as_attachment=True,
        download_name=f"{safe_title}.pptx",
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )


# ---------- WORKSHEET UPLOADS ----------

@app.route('/upload_excel_multi', methods=['POST'])
def upload_excel_multi():
    files = request.files.getlist('files')
    if not files:
        return jsonify({"error": "No worksheet files uploaded"}), 400

    try:
        raw_sheet_maps = []
        standard_sheet_maps = []

        for file in files:
            file_bytes = file.read()
            if not file_bytes:
                continue
            raw_sheet_maps.append(_read_excel_sheets_from_bytes(file_bytes, header=None))
            standard_sheet_maps.append(_read_excel_sheets_from_bytes(file_bytes, header=0))
    except Exception as e:
        return jsonify({"error": f"Could not read Excel files: {e}"}), 400

    if not raw_sheet_maps:
        return jsonify({"error": "Uploaded worksheet files were empty"}), 400

    merged_raw_sheets = _merge_sheet_maps(raw_sheet_maps)
    merged_standard_sheets = _merge_sheet_maps(standard_sheet_maps)

    species_list = _structured_life_history_workbook_to_species_list(merged_raw_sheets)
    if not species_list:
        species_list = _generic_workbook_to_species_list(merged_standard_sheets)

    if not species_list:
        return jsonify({"error": "Could not extract species and stage data from the uploaded worksheets"}), 400

    return jsonify(_dataset_response("Life History Worksheets", species_list))


@app.route('/upload_excel_single', methods=['POST'])
def upload_excel_single():
    if 'file' not in request.files:
        return jsonify({"error": "No worksheet file uploaded"}), 400

    file = request.files['file']

    try:
        file_bytes = file.read()
        if not file_bytes:
            return jsonify({"error": "Uploaded worksheet file was empty"}), 400
        raw_sheets = _read_excel_sheets_from_bytes(file_bytes, header=None)
        standard_sheets = _read_excel_sheets_from_bytes(file_bytes, header=0)
    except Exception as e:
        return jsonify({"error": f"Could not read Excel file: {e}"}), 400

    species_list = _structured_life_history_workbook_to_species_list(raw_sheets)
    if not species_list:
        first_sheet = next(iter(standard_sheets.values()), None)
        species_list = _simple_template_sheet_to_species_list(first_sheet)
    if not species_list:
        species_list = _generic_workbook_to_species_list(standard_sheets)

    if not species_list:
        return jsonify({"error": "Could not extract a life-history diagram from this worksheet"}), 400

    return jsonify(_dataset_response("Life History Worksheet", species_list))


# ---------- AI ENHANCEMENT ROUTE ----------

@app.route('/ai_enhance', methods=['POST'])
def ai_enhance():
    """
    Takes a single-species JSON (title, stages[])
    and asks an LLM to:
      - order stages logically,
      - merge duplicates,
      - compress / clean bullets.

    Returns: { title, stages[] } with same structure.

    IMPORTANT:
    - No new facts invented.
    - Keep bullets concise.
    """
    try:
        payload = request.get_json()
    except Exception:
        return jsonify({"error": "Invalid JSON"}), 400

    if not payload or 'stages' not in payload:
        return jsonify({"error": "Missing stages in payload"}), 400

    species_title = payload.get('title', 'Life history')
    stages = payload.get('stages', [])

    compact = {
        "title": species_title,
        "stages": [
            {"title": s.get("title", ""), "bullets": s.get("bullets", [])}
            for s in stages
        ]
    }

    system_prompt = (
        "You are an ecologist who designs clear life-history diagrams.\n"
        "TASK:\n"
        "- Reorder stages into a logical life-history sequence.\n"
        "- Merge near-duplicate stages when appropriate.\n"
        "- Preserve clear 'Duration:' bullets when present.\n"
        "- Preserve clear 'Timing:' / seasonal timing bullets when present.\n"
        "- Preserve one clear 'Lifespan:' value when present anywhere in the input.\n"
        "- For each stage, keep at most 5 concise bullets summarizing the information\n"
        "  (habitat, movement, food, reproduction, lifespan, etc.).\n"
        "- Do NOT invent any new biological facts; only rephrase or combine what is given.\n\n"
        "OUTPUT FORMAT (IMPORTANT):\n"
        "- Return ONLY valid JSON.\n"
        "- No explanation, no commentary, no markdown, no ``` fences.\n"
        "- Shape must be: {\"title\": string, \"stages\": [{\"title\": string, \"bullets\": [string, ...]}, ...]}\n"
    )

    user_prompt = (
        "Here is the current life-history data as JSON. "
        "Reorder and clean it as described, and return ONLY the JSON object.\n\n"
        f"{json.dumps(compact, ensure_ascii=False)}"
    )

    try:
        resp = client.responses.create(
            model="gpt-4.1-mini",
            input=[
                {"role": "system", "content": [{"type": "input_text", "text": system_prompt}]},
                {"role": "user", "content": [{"type": "input_text", "text": user_prompt}]},
            ],
        )

        text_out = resp.output[0].content[0].text.strip()

        # If model accidentally wraps with ```json fences, strip them
        if text_out.startswith("```"):
            lines = text_out.splitlines()
            lines = lines[1:]  # drop opening fence
            if lines and lines[-1].strip().startswith("```"):
                lines = lines[:-1]
            text_out = "\n".join(lines).strip()

        enhanced = json.loads(text_out)

    except Exception as e:
        return jsonify({"error": f"AI enhancement failed: {e}"}), 500

    if 'stages' not in enhanced or not isinstance(enhanced['stages'], list):
        return jsonify({"error": "AI returned invalid structure"}), 500

    return jsonify(enhanced)


if __name__ == '__main__':
    # Make sure you've set OPENAI_API_KEY in your environment
    app.run(debug=True, port=8000)
